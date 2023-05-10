from pptx import Presentation
from pptx.util import Inches
from os import walk, listdir, chdir, makedirs
from os.path import join as joinpath, dirname, abspath, splitext, exists



output_path = "/home/sysadmin/mounts/User-Primary--Sabari/Projects/Upwork/Automation/Python/PPTImbuseAudio/Output"
input_path = "/home/sysadmin/mounts/User-Primary--Sabari/Projects/Upwork/Automation/Python/PPTImbuseAudio/Input"

input_file = joinpath(input_path, "sample.pptx")
icon_path = joinpath(input_path , "sound-off-icon-40944.png")

input_files = listdir(input_path)
file_obj = Presentation(input_file)
print(file_obj.slide_width, file_obj.slide_height)
audio_map = {}

# # https://github.com/scanny/python-pptx/issues/323
# from pptx.opc.package import PartFactory
# from pptx.parts.media import MediaPart

# PartFactory.part_type_for.update(
#     {
#         'audio/mp3': MediaPart
#     }
# )

left = top = height = width = Inches(0.3)

def find_data(file_obj):
    
    for slide in file_obj.slides:
        hasPPT = False
        print(slide.name)
        for shape in slide.shapes:
            print(shape)
            if(hasattr(shape, "text")):
                pptid = shape.text
                if pptid.startswith("PPT ID:"):
                    pptid = pptid.replace("PPT ID:", "").strip()
                    hasPPT = True
                    audio_file = (audio_map[pptid])
                    break
        if hasPPT:
            slide.shapes.add_movie(audio_file, left,top,width,height, poster_frame_image=icon_path, mime_type="audio/mp3")
    file_obj.save(joinpath(output_path, "output.pptx"))



#         if hasPPT:
#             output_folder = joinpath(output_path, splitext(ifile)[0])
#             if not exists(output_folder):
#                 try:
#                     chdir(output_folder)
#                 except:
#                     makedirs(output_folder)
#                     pass
#             with open(joinpath(output_folder, pptid + ".txt"), "w") as f:
#                 f.write(slide.notes_slide.notes_text_frame.text.encode("utf-8").decode("utf-8"))

# def collect

for ifile in input_files:
    filename, ext = splitext(ifile)
    if ext == ".mp3":
        pptid = filename.split("-")[0]
        audio_map[pptid] = joinpath(input_path, ifile)

find_data(file_obj)

print("OUTPUT")
file_obj = Presentation(joinpath(output_path, "output.pptx"))
for slide in file_obj.slides:
    print(slide.name)
    for shape in slide.shapes:
        print(shape)