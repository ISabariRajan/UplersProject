from pptx import Presentation
from os import walk, listdir, chdir, makedirs
from os.path import join as joinpath, dirname, abspath, splitext, exists



output_path = "/home/sysadmin/mounts/User-Primary--Sabari/Projects/Upwork/Automation/Python/PPTXExtractor/Output"
input_path = "/home/sysadmin/mounts/User-Primary--Sabari/Projects/Upwork/Automation/Python/PPTXExtractor/Input"



def find_data(file_obj, ifile):
    for slide in file_obj.slides:
        if slide.notes_slide:
            # print(slide.notes_slide.notes_text_frame.text)
            hasPPT = False
            for shape in slide.shapes:
                if(hasattr(shape, "text")):
                    text = shape.text
                    if text.startswith("PPT ID:"):
                        text = text.replace("PPT ID:", "").strip()
                        hasPPT = True
                        print(text)
                        break

            if hasPPT:
                output_folder = joinpath(output_path, splitext(ifile)[0])
                if not exists(output_folder):
                    try:
                        chdir(output_folder)
                    except:
                        makedirs(output_folder)
                        pass
                with open(joinpath(output_folder, text + ".txt"), "w") as f:
                    f.write(slide.notes_slide.notes_text_frame.text.encode("utf-8").decode("utf-8"))


input_files = listdir(input_path)
for ifile in input_files:
    if splitext(ifile)[1] == ".pptx":
        file_obj = Presentation(joinpath(input_path, ifile))
        find_data(file_obj, ifile)
    
        

    print("\n")